# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PALETTE
BG_DARK     = RGBColor(0x0D, 0x0D, 0x0D)
GOLD        = RGBColor(0xF5, 0xC5, 0x18)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN       = RGBColor(0x86, 0xEF, 0xAC)
RED_SOFT    = RGBColor(0xF8, 0x71, 0x71)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── HELPERS
def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_tb(slide, text, l, t, w, h, size=14, bold=False,
           color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def gold_line(slide, t):
    add_rect(slide, 0, t, W, Inches(0.04), GOLD)

def bullets(slide, items, l, t, w, size=14, color=LIGHT_GRAY):
    lh = Inches(0.4)
    for item in items:
        prefix = "  \u2022  " if not item.startswith("    ") else "      \u25b8  "
        add_tb(slide, prefix + item.strip(), l, t, w, lh, size=size, color=color)
        t += lh
    return t

def code_block(slide, code, l, t, w):
    lines = code.strip().split("\n")
    h = Inches(0.3) * len(lines) + Inches(0.2)
    add_rect(slide, l, t, w, h, DARK_GRAY)
    add_tb(slide, code.strip(), l+Inches(0.1), t+Inches(0.08),
           w-Inches(0.2), h, size=11, color=GOLD)
    return t + h

def content_slide(title, speaker, fn):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), GOLD)
    add_rect(slide, 0, Inches(0.08), Inches(3.0), Inches(0.42), RGBColor(0x1E,0x1E,0x1E))
    add_tb(slide, "  " + speaker, Inches(0.05), Inches(0.08), Inches(2.9), Inches(0.42),
           size=11, color=GOLD)
    add_tb(slide, title, Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.75),
           size=26, bold=True, color=WHITE)
    gold_line(slide, Inches(1.42))
    fn(slide)
    return slide

# ═══════════════════════════════════════
# TITLE SLIDE
# ═══════════════════════════════════════
def title_slide():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.12), GOLD)
    add_rect(slide, 0, H-Inches(0.12), W, Inches(0.12), GOLD)
    add_tb(slide, "X-AURUM", Inches(1), Inches(1.1), Inches(11), Inches(1.4),
           size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_tb(slide, "H\u1ec7 Th\u1ed1ng D\u1ef1 B\u00e1o Gi\u00e1 V\u00e0ng Theo Th\u1eddi Gian Th\u1ef1c",
           Inches(1), Inches(2.65), Inches(11), Inches(0.75),
           size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(slide, "Tri\u1ec3n khai Web API t\u00edch h\u1ee3p Tr\u00ed tu\u1ec7 Nh\u00e2n t\u1ea1o & D\u1eef li\u1ec7u Th\u1ecb tr\u01b0\u1eddng To\u00e0n c\u1ea7u",
           Inches(1), Inches(3.35), Inches(11), Inches(0.6),
           size=15, color=LIGHT_GRAY, italic=True, align=PP_ALIGN.CENTER)
    gold_line(slide, Inches(4.1))
    add_tb(slide, "Nh\u00f3m 3  |  Xu\u00e2n H\u01b0\u01a1ng  |  Qu\u1ed1c \u0110\u1ea1t  |  V\u0103n Nguy\u1ec5n  |  Th\u00e0nh Trung",
           Inches(1), Inches(4.3), Inches(11), Inches(0.55),
           size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_tb(slide, "nttspace.online", Inches(1), Inches(5.0), Inches(11), Inches(0.55),
           size=14, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

title_slide()

# ═══════════════════════════════════════
# SLIDES 1-3: Xuan Huong — API Provisioning
# ═══════════════════════════════════════
def s1(slide):
    add_tb(slide, "V\u1ea5n \u0111\u1ec1 \u0111\u1eb7t ra", Inches(0.5), Inches(1.58), Inches(6), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "M\u00f4 h\u00ecnh AI (ML.NET) ch\u1ec9 ch\u1ea1y \u0111\u01b0\u1ee3c tr\u00ean Console n\u1ed9i b\u1ed9.",
        "Kh\u00f4ng th\u1ec3 chia s\u1ebb ho\u1eb7c t\u00e1i s\u1eed d\u1ee5ng t\u1eeb b\u00ean ngo\u00e0i.",
    ], Inches(0.5), Inches(2.0), Inches(5.8))
    add_tb(slide, "Gi\u1ea3i ph\u00e1p: ASP.NET Core Minimal API (.NET 10)", Inches(0.5), Inches(2.9), Inches(6), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "B\u1ecdc thu\u1eadt to\u00e1n AI th\u00e0nh Web Service chu\u1ea9n HTTP.",
        "Lo\u1ea1i b\u1ecf Controllers truy\u1ec1n th\u1ed1ng \u2014 ki\u1ebfn tr\u00fac Microservices tinh g\u1ecdn nh\u1ea5t.",
    ], Inches(0.5), Inches(3.32), Inches(5.8))
    add_tb(slide, "Clean Architecture", Inches(7), Inches(1.58), Inches(5.8), Inches(0.38),
           size=15, bold=True, color=GOLD)
    code_block(slide,
        "GoldPrice_API/\n"
        "  Endpoints/  -> Dinh tuyen API Routes\n"
        "  Models/     -> Input / Output Schema\n"
        "  Extensions/ -> DI, RateLimit, ML Pool",
        Inches(7), Inches(2.0), Inches(5.8))

content_slide("Slide 1 \u00b7 B\u00e0i To\u00e1n & Ki\u1ebfn Tr\u00fac C\u1ed1t L\u00f5i", "Xu\u00e2n H\u01b0\u01a1ng", s1)

def s2(slide):
    cols_x = [Inches(0.5), Inches(3.8), Inches(7.6)]
    col_w  = [Inches(3.3), Inches(1.8), Inches(5.0)]
    hdrs   = ["Endpoint", "Method", "M\u00f4 t\u1ea3"]
    for i, (x, hdr) in enumerate(zip(cols_x, hdrs)):
        add_rect(slide, x, Inches(1.58), col_w[i], Inches(0.4), GOLD)
        add_tb(slide, hdr, x+Inches(0.05), Inches(1.58), col_w[i], Inches(0.4),
               size=13, bold=True, color=BG_DARK)
    rows = [
        ["/api/v1/predictions",          "POST", "Client t\u1ef1 nh\u1eadp d\u1eef li\u1ec7u OHLCV"],
        ["/api/v1/predictions/realtime", "GET",  "Bot t\u1ef1 k\u00e9o data \u2014 kh\u00f4ng c\u1ea7n input"],
    ]
    for r, row in enumerate(rows):
        bg = RGBColor(0x1A,0x1A,0x1A) if r % 2 == 0 else RGBColor(0x22,0x22,0x22)
        for i, (x, cell) in enumerate(zip(cols_x, row)):
            add_rect(slide, x, Inches(1.98+r*0.45), col_w[i], Inches(0.45), bg)
            clr = GOLD if i == 1 else LIGHT_GRAY
            add_tb(slide, cell, x+Inches(0.05), Inches(1.98+r*0.45), col_w[i], Inches(0.45),
                   size=12, color=clr)
    add_tb(slide, "Chu\u1ea9n h\u00f3a Response RESTful (JSON Wrapper)", Inches(0.5), Inches(3.1), Inches(6.5), Inches(0.38),
           size=14, bold=True, color=GOLD)
    code_block(slide,
        '{\n  "success": true,\n  "data": {\n'
        '    "predictedClose": 2355.8,\n'
        '    "predictedCloseVnd": 120000000,\n'
        '    "liveWorldPrice": 2341.2\n  }\n}',
        Inches(0.5), Inches(3.52), Inches(6.5))
    bullets(slide, [
        "M\u1ecdi Client \u0111\u1ec1u parse an to\u00e0n.",
        "D\u1ec5 m\u1edf r\u1ed9ng th\u00eam tr\u01b0\u1eddng m\u1edbi.",
        "Ph\u00e2n bi\u1ec7t r\u00f5 lu\u1ed3ng Manual vs Auto.",
    ], Inches(7.3), Inches(3.12), Inches(5.6), size=14)

content_slide("Slide 2 \u00b7 Thi\u1ebft K\u1ebf Hai Lu\u1ed3ng Endpoint", "Xu\u00e2n H\u01b0\u01a1ng", s2)

def s3(slide):
    add_tb(slide, "T\u1ea1i sao c\u1ea7n Swagger / OpenAPI?", Inches(0.5), Inches(1.58), Inches(12), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "API \u0111\u01b0a ra ngo\u00e0i m\u00e0 kh\u00f4ng c\u00f3 m\u00f4 t\u1ea3 \u2192 B\u00ean th\u1ee9 3 kh\u00f4ng bi\u1ebft c\u00e1ch d\u00f9ng.",
        "OpenAPI l\u00e0 chu\u1ea9n m\u00f4 t\u1ea3 API \u0111\u01b0\u1ee3c c\u1ea3 ng\u00e0nh c\u00f4ng nghi\u1ec7p s\u1eed d\u1ee5ng.",
        "Swagger UI t\u1ef1 \u0111\u1ed9ng \u0111\u1ecdc code C# \u2192 Sinh trang t\u00e0i li\u1ec7u t\u01b0\u01a1ng t\u00e1c \u0111\u1ea7y \u0111\u1ee7.",
    ], Inches(0.5), Inches(2.0), Inches(12.3))
    add_tb(slide, "T\u00ednh n\u0103ng \u0111\u00e3 t\u00edch h\u1ee3p", Inches(0.5), Inches(3.25), Inches(12), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "Hi\u1ec3n th\u1ecb \u0111\u1ea7y \u0111\u1ee7 2 endpoints, m\u00f4 t\u1ea3 tham s\u1ed1 v\u00e0 response schema.",
        "Cho ph\u00e9p Execute request ngay tr\u00ean giao di\u1ec7n Web \u2014 kh\u00f4ng c\u1ea7n Postman.",
        "L\u00e0 c\u00f4ng c\u1ee5 Demo tr\u1ef1c ti\u1ebfp trong bu\u1ed5i b\u00e1o c\u00e1o h\u00f4m nay.",
    ], Inches(0.5), Inches(3.67), Inches(12.3))
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.72), RGBColor(0x1B,0x2A,0x1B))
    add_tb(slide, "Swagger ch\u00ednh l\u00e0 c\u1ed5ng giao ti\u1ebfp gi\u1eefa nh\u00f3m v\u00e0 th\u1ebf gi\u1edbi b\u00ean ngo\u00e0i.",
           Inches(0.7), Inches(5.45), Inches(12), Inches(0.62),
           size=14, color=GREEN, italic=True, align=PP_ALIGN.CENTER)

content_slide("Slide 3 \u00b7 T\u00e0i Li\u1ec7u API T\u1ef1 \u0110\u1ed9ng V\u1edbi Swagger (OpenAPI)", "Xu\u00e2n H\u01b0\u01a1ng", s3)

# ═══════════════════════════════════════
# SLIDES 4-6: Quoc Dat — 3rd-Party APIs
# ═══════════════════════════════════════
def s4(slide):
    add_tb(slide, "B\u00e0i thi 3: S\u1eed d\u1ee5ng API t\u1eeb ngu\u1ed3n b\u00ean ngo\u00e0i (3rd-Party APIs)",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.42),
           size=16, bold=True, color=GOLD)
    add_tb(slide, "Server kh\u00f4ng ch\u1edd User nh\u1eadp li\u1ec7u \u2014 t\u1ef1 \u0111\u1ed9ng g\u1ecdi ra 2 ngu\u1ed3n d\u1eef li\u1ec7u th\u1ef1c t\u1ebf:",
           Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.35),
           size=13, color=LIGHT_GRAY, italic=True)
    # Card Binance
    add_rect(slide, Inches(0.5), Inches(2.5), Inches(6.0), Inches(3.2), RGBColor(0x12,0x12,0x28))
    add_rect(slide, Inches(0.5), Inches(2.5), Inches(6.0), Inches(0.42), GOLD)
    add_tb(slide, "1. Binance Public API", Inches(0.6), Inches(2.5), Inches(5.8), Inches(0.42),
           size=14, bold=True, color=BG_DARK)
    add_tb(slide, "M\u1ee5c \u0111\u00edch: L\u1ea5y gi\u00e1 V\u00e0ng Spot theo phi\u00ean giao d\u1ecbch",
           Inches(0.65), Inches(2.97), Inches(5.7), Inches(0.35), size=12, color=LIGHT_GRAY, italic=True)
    add_tb(slide, "HTTP Method & URL:", Inches(0.65), Inches(3.36), Inches(5.7), Inches(0.3),
           size=11, bold=True, color=GOLD)
    code_block(slide,
        "GET api.binance.com/api/v3/klines\n"
        "  ?symbol=PAXGUSDT&interval=1d&limit=1",
        Inches(0.65), Inches(3.7), Inches(5.7))
    add_tb(slide, "Tr\u00edch xu\u1ea5t: Open, High, Low, Volume (t\u1eeb m\u1ea3ng Klines JSON)",
           Inches(0.65), Inches(4.45), Inches(5.7), Inches(0.35), size=12, color=GREEN)
    # Card Exchange Rates
    add_rect(slide, Inches(7.0), Inches(2.5), Inches(6.0), Inches(3.2), RGBColor(0x12,0x12,0x28))
    add_rect(slide, Inches(7.0), Inches(2.5), Inches(6.0), Inches(0.42), GOLD)
    add_tb(slide, "2. Open Exchange Rates API", Inches(7.1), Inches(2.5), Inches(5.8), Inches(0.42),
           size=14, bold=True, color=BG_DARK)
    add_tb(slide, "M\u1ee5c \u0111\u00edch: L\u1ea5y t\u1ef7 gi\u00e1 USD \u2192 VND hi\u1ec7n t\u1ea1i",
           Inches(7.15), Inches(2.97), Inches(5.7), Inches(0.35), size=12, color=LIGHT_GRAY, italic=True)
    add_tb(slide, "HTTP Method & URL:", Inches(7.15), Inches(3.36), Inches(5.7), Inches(0.3),
           size=11, bold=True, color=GOLD)
    code_block(slide, "GET open.er-api.com/v6/latest/USD", Inches(7.15), Inches(3.7), Inches(5.7))
    add_tb(slide, "Tr\u00edch xu\u1ea5t: rates.VND (t\u1eeb object JSON ph\u1ea3n h\u1ed3i)",
           Inches(7.15), Inches(4.45), Inches(5.7), Inches(0.35), size=12, color=GREEN)
    add_rect(slide, Inches(0.5), Inches(5.85), Inches(12.5), Inches(0.62), RGBColor(0x1B,0x2A,0x1B))
    add_tb(slide, "K\u1ebft qu\u1ea3: PredictedClose (USD)  x  t\u1ef7 gi\u00e1 VND  =  Gi\u00e1 V\u00e0ng VND hi\u1ec7n t\u1ea1i",
           Inches(0.7), Inches(5.9), Inches(12.2), Inches(0.52),
           size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

content_slide("Slide 4 \u00b7 T\u00edch H\u1ee3p API B\u00ean Th\u1ee9 3 (Binance & ExchangeRate)", "Qu\u1ed1c \u0110\u1ea1t", s4)

def s5(slide):
    add_tb(slide, "Lu\u1ed3ng x\u1eed l\u00fd b\u00ean trong /realtime endpoint (b\u1ea5t \u0111\u1ed3ng b\u1ed9):",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    code_block(slide,
        "1. Server gui HTTP GET -> Binance API\n"
        "   -> Nhan JSON Klines -> Trich Open, High, Low, Volume\n\n"
        "2. Server gui HTTP GET -> Exchange Rates API\n"
        "   -> Doc node 'VND' -> Nhan ty gia\n\n"
        "3. Gop 4 gia tri -> Model AI tinh toan\n"
        "   -> AI tra ve PredictedClose (USD)\n\n"
        "4. PredictedClose x VND_Rate = Gia VND\n"
        "   -> Tra ve JSON cho Client",
        Inches(0.5), Inches(2.0), Inches(7.0))
    add_tb(slide, "T\u1ea1i sao d\u00f9ng async/await?", Inches(7.8), Inches(1.58), Inches(5.1), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "Server kh\u00f4ng b\u1ecb '\u0111\u1ee9ng' khi ch\u1edd m\u1ea1ng.",
        "X\u1eed l\u00fd nhi\u1ec1u request \u0111\u1ed3ng th\u1eddi.",
        "T\u0103ng throughput t\u1ed5ng th\u1ec3.",
        "",
        "K\u1ef9 thu\u1eadt: Task<IResult>",
        "JSON: System.Text.Json DOM tree",
        "\u2192 Ch\u1ec9 \u0111\u1ecdc \u0111\u00fang node c\u1ea7n thi\u1ebft,",
        "   kh\u00f4ng map to\u00e0n b\u1ed9 object.",
    ], Inches(7.8), Inches(2.0), Inches(5.1), size=13)

content_slide("Slide 5 \u00b7 X\u1eed L\u00fd B\u1ea5t \u0110\u1ed3ng B\u1ed9 & Ph\u00e2n T\u00edch JSON", "Qu\u1ed1c \u0110\u1ea1t", s5)

def s6(slide):
    add_tb(slide, "V\u1ea5n \u0111\u1ec1 khi d\u00f9ng new HttpClient() tr\u1ef1c ti\u1ebfp:",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "M\u1ed7i Request t\u1ea1o ra m\u1ed9t k\u1ebft n\u1ed1i Socket m\u1edbi \u2192 C\u1ea1n ki\u1ec7t Port (Socket Exhaustion).",
        "\u1ea2nh h\u01b0\u1edfng nghi\u00eam tr\u1ecdng khi c\u00f3 h\u00e0ng tr\u0103m ng\u01b0\u1eddi d\u00f9ng g\u1ecdi API \u0111\u1ed3ng th\u1eddi.",
    ], Inches(0.5), Inches(2.0), Inches(12.3))
    add_tb(slide, "Gi\u1ea3i ph\u00e1p: IHttpClientFactory",
           Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "\u0110\u0103ng k\u00fd v\u00e0o Dependency Injection Container \u2014 qu\u1ea3n l\u00fd v\u00f2ng \u0111\u1eddi b\u1ed9 g\u1ecdi HTTP.",
        "T\u00e1i s\u1eed d\u1ee5ng k\u1ebft n\u1ed1i Socket hi\u1ec7u qu\u1ea3 \u2192 Kh\u00f4ng bao gi\u1edd b\u1ecb c\u1ea1n ki\u1ec7t Port.",
    ], Inches(0.5), Inches(3.27), Inches(12.3))
    add_tb(slide, "C\u1ea5u h\u00ecnh \u2014 ch\u1ec9 1 d\u00f2ng trong ServiceExtensions.cs:",
           Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.35),
           size=13, bold=True, color=LIGHT_GRAY)
    code_block(slide, "builder.Services.AddHttpClient();",
               Inches(0.5), Inches(4.44), Inches(12.3))
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.72), RGBColor(0x1A,0x1A,0x2E))
    add_tb(slide, "Thay \u0111\u1ed5i nh\u1ecf trong DI Container, ng\u0103n ch\u1eb7n to\u00e0n b\u1ed9 nguy c\u01a1 Port Exhaustion khi scale.",
           Inches(0.7), Inches(5.45), Inches(12), Inches(0.62),
           size=14, color=ACCENT_BLUE, italic=True, align=PP_ALIGN.CENTER)

content_slide("Slide 6 \u00b7 IHttpClientFactory - G\u1ecdi API Chuy\u00ean Nghi\u1ec7p", "Qu\u1ed1c \u0110\u1ea1t", s6)

# ═══════════════════════════════════════
# SLIDES 7-9: Van Nguyen — Memory & Security
# ═══════════════════════════════════════
def s7(slide):
    add_tb(slide, "V\u1ea5n \u0111\u1ec1 \u0111\u1eb7c th\u00f9 khi Host Machine Learning l\u00e0m REST API:",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "File GoldModel.zip n\u1eb7ng ~43MB \u2014 ch\u1ee9a to\u00e0n b\u1ed9 c\u1ea5u tr\u00fac th\u1ea7n kinh nh\u00e2n t\u1ea1o.",
        "N\u1ebfu n\u1ea1p l\u1ea1i model m\u1ed7i l\u1ea7n c\u00f3 Request \u2192 T\u1ed1n h\u00e0ng gi\u00e2y x\u1eed l\u00fd m\u1ed7i l\u01b0\u1ee3t g\u1ecdi.",
        "H\u00e0ng tr\u0103m ng\u01b0\u1eddi d\u00f9ng \u0111\u1ed3ng th\u1eddi \u2192 Memory Leak \u2192 Server Crash.",
    ], Inches(0.5), Inches(2.0), Inches(12.3))
    add_tb(slide, "K\u1ecbch b\u1ea3n nguy hi\u1ec3m:", Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=RED_SOFT)
    code_block(slide,
        "// SAI -- nap lai model moi request\n"
        "app.MapPost('/predict', (input) => {\n"
        "    var ctx = new MLContext();\n"
        "    var model = ctx.Model.Load('GoldModel.zip', ...);\n"
        "    // -> RAM tang vo han theo so luong request\n"
        "});",
        Inches(0.5), Inches(3.62), Inches(12.3))

content_slide("Slide 7 \u00b7 B\u00e0i To\u00e1n Hi\u1ec7u N\u0103ng - Hosting ML Model", "V\u0103n Nguy\u1ec5n", s7)

def s8(slide):
    add_tb(slide, "Gi\u1ea3i ph\u00e1p: PredictionEnginePool (Object Pooling)",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    code_block(slide,
        "// DUNG -- dang ky 1 lan duy nhat luc Server khoi dong\n"
        "builder.Services.AddPredictionEnginePool<ModelInput, ModelOutput>(\n"
        "    modelName: 'GoldModel',\n"
        "    modelBuilder: ctx => ctx.Model.Load('GoldModel.zip', ...)\n"
        ");",
        Inches(0.5), Inches(2.0), Inches(12.3))
    add_tb(slide, "C\u01a1 ch\u1ebf ho\u1ea1t \u0111\u1ed9ng:", Inches(0.5), Inches(3.28), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "Server start \u2192 N\u1ea1p model v\u00e0o RAM \u0111\u00fang 1 l\u1ea7n.",
        "Request \u0111\u1ebfn \u2192 M\u01b0\u1ee3n Engine \u2192 D\u1ef1 b\u00e1o \u2192 Tr\u1ea3 l\u1ea1i Pool.",
    ], Inches(0.5), Inches(3.7), Inches(6.2), size=14)
    bullets(slide, [
        "Thread-Safe: Ng\u00e0n requests song song, kh\u00f4ng xung \u0111\u1ed9t.",
        "Zero Memory Leak: Pool t\u1ef1 qu\u1ea3n l\u00fd v\u00f2ng \u0111\u1eddi object.",
    ], Inches(6.8), Inches(3.7), Inches(6.1), size=14)

content_slide("Slide 8 \u00b7 PredictionEnginePool - Qu\u1ea3n Tr\u1ecb RAM", "V\u0103n Nguy\u1ec5n", s8)

def s9(slide):
    add_tb(slide, "Nguy c\u01a1: API c\u00f4ng khai \u2192 D\u1ec5 b\u1ecb DDoS (t\u1ea5n c\u00f4ng l\u00e0m ngh\u1ebdn Server)",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    add_tb(slide, "Gi\u1ea3i ph\u00e1p: Fixed Window Rate Limiting Middleware",
           Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    add_rect(slide, Inches(0.5), Inches(2.55), Inches(5.7), Inches(2.2), RGBColor(0x1A,0x2A,0x1A))
    add_tb(slide, "C\u1ea5u h\u00ecnh hi\u1ec7n t\u1ea1i:", Inches(0.65), Inches(2.62), Inches(5.3), Inches(0.35),
           size=13, bold=True, color=GREEN)
    add_tb(slide, "  Permit Limit: 5 Request\n  Window: 10 gi\u00e2y / IP\n  Queue: 2 Request \u0111\u01b0\u1ee3c x\u1ebfp h\u00e0ng",
           Inches(0.65), Inches(3.02), Inches(5.3), Inches(1.5), size=14, color=WHITE)
    add_rect(slide, Inches(6.5), Inches(2.55), Inches(6.3), Inches(2.2), RGBColor(0x2A,0x1A,0x1A))
    add_tb(slide, "K\u1ebft qu\u1ea3 khi v\u01b0\u1ee3t gi\u1edbi h\u1ea1n:", Inches(6.65), Inches(2.62), Inches(5.9), Inches(0.35),
           size=13, bold=True, color=RED_SOFT)
    add_tb(slide, "  HTTP 429 Too Many Requests\n\n  B\u1ecb ch\u1eb7n t\u1ea1i t\u1ea7ng HTTP Pipeline\n  \u2192 Ch\u01b0a ch\u1ea1m v\u00e0o AI Model\n  \u2192 CPU & RAM ho\u00e0n to\u00e0n an to\u00e0n",
           Inches(6.65), Inches(3.02), Inches(5.9), Inches(1.5), size=13, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.72), RGBColor(0x1B,0x1B,0x2A))
    add_tb(slide, "Middleware ch\u1eb7n t\u1eeb ngo\u00e0i v\u00e0o \u2014 l\u1edbp khi\u00ean th\u00e9p b\u1ea3o v\u1ec7 to\u00e0n b\u1ed9 t\u00e0i nguy\u00ean ph\u00eda trong.",
           Inches(0.7), Inches(5.45), Inches(12), Inches(0.62),
           size=14, color=ACCENT_BLUE, italic=True, align=PP_ALIGN.CENTER)

content_slide("Slide 9 \u00b7 Rate Limiting - T\u01b0\u1eddng L\u1eeda Ch\u1ed1ng T\u1ea5n C\u00f4ng", "V\u0103n Nguy\u1ec5n", s9)

# ═══════════════════════════════════════
# SLIDES 10-12: Thanh Trung — DevOps & Demo
# ═══════════════════════════════════════
def s10(slide):
    add_tb(slide, "V\u1ea5n \u0111\u1ec1: 'Code ch\u1ea1y m\u00e1y em, sao Server l\u1ea1i l\u1ed7i?' \u2192 Xung \u0111\u1ed9t m\u00f4i tr\u01b0\u1eddng .NET",
           Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.38),
           size=15, bold=True, color=GOLD)
    add_tb(slide, "Gi\u1ea3i ph\u00e1p: Docker Multi-stage Build",
           Inches(0.5), Inches(2.05), Inches(6.5), Inches(0.38),
           size=15, bold=True, color=GOLD)
    code_block(slide,
        "# Stage 1: Build (SDK ~800MB)\n"
        "FROM mcr.microsoft.com/dotnet/sdk:10.0 AS build\n"
        "RUN dotnet publish -c Release -o /app/publish\n\n"
        "# Stage 2: Runtime (~200MB)\n"
        "FROM mcr.microsoft.com/dotnet/aspnet:10.0 AS final\n"
        "COPY --from=build /app/publish .\n"
        'ENTRYPOINT ["dotnet", "GoldPrice_API.dll"]',
        Inches(0.5), Inches(2.47), Inches(6.5))
    add_tb(slide, "K\u1ebft qu\u1ea3:", Inches(7.2), Inches(2.05), Inches(5.7), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "Image production ch\u1ec9 ~200MB.",
        "Kh\u00f4ng ch\u1ee9a SDK th\u1eeba.",
        "Ch\u1ea1y gi\u1ed1ng h\u1ec7t nhau tr\u00ean m\u1ecdi m\u00e1y c\u00f3 Docker.",
        "",
        "docker compose up -d",
        "\u2192 To\u00e0n b\u1ed9 h\u1ec7 th\u1ed1ng kh\u1edfi \u0111\u1ed9ng",
        "   ch\u1ec9 v\u1edbi 1 l\u1ec7nh duy nh\u1ea5t.",
    ], Inches(7.2), Inches(2.47), Inches(5.7), size=13)

content_slide("Slide 10 \u00b7 Docker - \u0110\u00f3ng G\u00f3i & Chu\u1ea9n H\u00f3a M\u00f4i Tr\u01b0\u1eddng", "Th\u00e0nh Trung", s10)

def s11(slide):
    add_tb(slide, "Quy tr\u00ecnh CI/CD t\u1ef1 \u0111\u1ed9ng (GitHub Actions):",
           Inches(0.5), Inches(1.58), Inches(7), Inches(0.38), size=15, bold=True, color=GOLD)
    steps = [
        "git push l\u00ean GitHub",
        "GitHub Actions k\u00edch ho\u1ea1t",
        "Build Docker Image",
        "Push l\u00ean GHCR (GitHub Container Registry)",
        "Server k\u00e9o image m\u1edbi & Restart",
    ]
    t = Inches(2.05)
    for i, step in enumerate(steps):
        clr = GOLD if i == 0 else (GREEN if i == 4 else LIGHT_GRAY)
        add_rect(slide, Inches(0.7), t, Inches(5.8), Inches(0.48), RGBColor(0x1E,0x1E,0x1E))
        prefix = "  >> " if i > 0 else "  [>>] "
        add_tb(slide, prefix + step, Inches(0.85), t+Inches(0.05), Inches(5.5), Inches(0.38),
               size=13, color=clr, bold=(i == 0 or i == 4))
        t += Inches(0.53)
    add_tb(slide, "Nginx Reverse Proxy:", Inches(7.2), Inches(1.58), Inches(5.7), Inches(0.38),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "Ti\u1ebfp nh\u1eadn traffic t\u1eeb nttspace.online",
        "(Port 80 / 443 HTTPS).",
        "",
        "Chuy\u1ec3n ti\u1ebfp v\u00e0o Docker Container",
        "\u0111ang l\u1eafng nghe n\u1ed9i b\u1ed9.",
        "",
        "K\u1ebft qu\u1ea3:",
    ], Inches(7.2), Inches(2.05), Inches(5.7), size=13)
    add_rect(slide, Inches(7.2), Inches(5.05), Inches(5.7), Inches(0.58), GOLD)
    add_tb(slide, "nttspace.online  --  LIVE!", Inches(7.3), Inches(5.1), Inches(5.5), Inches(0.48),
           size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

content_slide("Slide 11 \u00b7 CI/CD & Tri\u1ec3n Khai Tr\u00ean nttspace.online", "Th\u00e0nh Trung", s11)

def s12(slide):
    add_tb(slide, "Demo 1 \u2014 GET /api/v1/predictions/realtime",
           Inches(0.5), Inches(1.55), Inches(6.2), Inches(0.42),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "1. M\u1edf Swagger UI tr\u00ean nttspace.online",
        "2. Ch\u1ecdn GET /realtime \u2192 Execute",
        "3. Server t\u1ef1 g\u1ecdi Binance + Exchange Rates",
        "4. AI t\u00ednh to\u00e1n \u2192 Tr\u1ea3 JSON:",
    ], Inches(0.5), Inches(2.02), Inches(6.2), size=13)
    code_block(slide,
        '"predictedClose": 2355.8\n"predictedCloseVnd": 120,000,000 VND',
        Inches(0.5), Inches(3.75), Inches(6.2))
    add_tb(slide, "Demo 2 \u2014 Ki\u1ec3m tra Rate Limiting",
           Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.42),
           size=15, bold=True, color=GOLD)
    bullets(slide, [
        "1. Ch\u1ecdn POST /predictions \u2192 Execute li\u00ean t\u1ee5c",
        "2. L\u1ea7n 1-5: Response 200 OK",
        "3. L\u1ea7n 6-7: Response 429 [!]",
        "   Too Many Requests",
        "\u2192 T\u01b0\u1eddng l\u1eeda ho\u1ea1t \u0111\u1ed9ng ho\u00e0n h\u1ea3o!",
    ], Inches(7.0), Inches(2.02), Inches(6.0), size=13, color=LIGHT_GRAY)
    add_rect(slide, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.82), RGBColor(0x1B,0x2A,0x1B))
    add_tb(slide, "To\u00e0n b\u1ed9 h\u1ec7 th\u1ed1ng \u0111ang ch\u1ea1y LIVE t\u1ea1i  nttspace.online",
           Inches(0.7), Inches(5.4), Inches(12), Inches(0.72),
           size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

content_slide("Slide 12 \u00b7 Live Demo Tr\u00ean Swagger", "Th\u00e0nh Trung", s12)

# ═══════════════════════════════════════
# FINAL SLIDE
# ═══════════════════════════════════════
def final_slide():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.12), GOLD)
    add_rect(slide, 0, H-Inches(0.12), W, Inches(0.12), GOLD)
    add_tb(slide, "C\u1ea2M \u01a0N TH\u1ea6Y C\u00d4 & H\u1ed8I \u0110\u1ed4NG!",
           Inches(1), Inches(1.0), Inches(11), Inches(1.2),
           size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_tb(slide, "Nh\u00f3m 3 xin l\u1eafng nghe c\u00e2u h\u1ecfi ph\u1ea3n bi\u1ec7n.",
           Inches(1), Inches(2.3), Inches(11), Inches(0.58),
           size=18, color=LIGHT_GRAY, italic=True, align=PP_ALIGN.CENTER)
    gold_line(slide, Inches(3.08))
    members = [
        ("Xu\u00e2n H\u01b0\u01a1ng",  "Ki\u1ebfn tr\u00fac API & Swagger"),
        ("Qu\u1ed1c \u0110\u1ea1t",    "T\u00edch h\u1ee3p 3rd-Party APIs"),
        ("V\u0103n Nguy\u1ec5n",  "Hi\u1ec7u n\u0103ng & B\u1ea3o m\u1eadt"),
        ("Th\u00e0nh Trung", "DevOps & Live Demo"),
    ]
    col_w = Inches(3.1)
    for i, (name, role) in enumerate(members):
        x = Inches(0.35) + i * Inches(3.25)
        add_rect(slide, x, Inches(3.28), col_w, Inches(1.18), RGBColor(0x1A,0x1A,0x1A))
        add_rect(slide, x, Inches(3.28), col_w, Inches(0.08), GOLD)
        add_tb(slide, name, x+Inches(0.1), Inches(3.38), col_w-Inches(0.2), Inches(0.48),
               size=15, bold=True, color=WHITE)
        add_tb(slide, role, x+Inches(0.1), Inches(3.88), col_w-Inches(0.2), Inches(0.52),
               size=12, color=LIGHT_GRAY, italic=True)
    add_tb(slide, "nttspace.online",
           Inches(1), Inches(5.08), Inches(11), Inches(0.58),
           size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

final_slide()

output = r"D:\Coding Space\Project\Group3\X-AURUM_Presentation.pptx"
prs.save(output)
print("[OK] PPTX saved: " + output + " (14 slides)")
